from pptx import Presentation
from deck2pptx.height_estimator import extract_template_metrics

def inspect_calibration(pptx_path: str):
    print(f"Inspecting Calibration Data for: {pptx_path}")
    try:
        prs = Presentation(pptx_path)
    except Exception as e:
        print(f"Error loading {pptx_path}: {e}")
        return

    if not prs.slides:
        print("No slides found in presentation.")
        return

    slide = prs.slides[0]
    print("\n--- Calibration data from Slide 1 ---")
    
    EMU_TO_INCH = 1.0 / 914400.0

    for i, shape in enumerate(slide.shapes):
        if not shape.has_text_frame:
            continue

        text = shape.text
        lines_raw = text.split('\n')
        lines_count = text.count('\n') + text.count('\x0b') + 1

        font_size_pt = None
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                if run.font.size:
                    font_size_pt = run.font.size.pt
                    break
            if font_size_pt: break
        
        h_in = shape.height * EMU_TO_INCH if shape.height else 0
        w_in = shape.width * EMU_TO_INCH if shape.width else 0
        
        print(f"\nShape {i}: {shape.name} (Font: {font_size_pt}pt)")
        print(f"  Dimensions: Width={w_in:.3f}in, Height={h_in:.3f}in")
        print(f"  Total Lines: {lines_count}")
        
        if lines_count > 0:
            h_per_line_in = h_in / lines_count
            print(f"  Avg Height per line: {h_per_line_in:.3f} in")
            
            # Print line lengths
            for j, line in enumerate(lines_raw):
                print(f"    Line {j+1}: length={len(line)} chars | '{line[:30]}...'")

            # Try to calculate CPI (Characters per inch) for the first line
            first_line = lines_raw[0]
            if w_in > 0 and font_size_pt:
                cpi = len(first_line) / w_in
                print(f"  Calculated CPI: {cpi:.2f} chars/inch")

    # Call actual implementation
    calibrated_metrics, title_metrics = extract_template_metrics(slide)

    if calibrated_metrics:
        print("\n--- Extracted Calibrated Metrics ---")
        for fs, metrics in calibrated_metrics.items():
            print(f"Font size {fs}pt -> height: {metrics['height']} EMU, CPI: {metrics['cpi']:.2f}")
    else:
        print("\nNo valid multi-line calibration data found for AiDeckCore standard.")
