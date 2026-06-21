from pptx.util import Inches
from ..render_context import SlideContext
from ..height_estimator import ELEMENT_GAP

def render(element, ctx: SlideContext, x, y, w, h) -> float:
    from . import render_element
    from ..height_estimator import get_adjusted_height
    
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    num_panels = len(element.panels)
    if num_panels == 0:
        return y
    target_x = ph.left if ph else x
    target_y = ph.top if ph else y
    target_w = ph.width if ph else w
    target_h = ph.height if ph else h
    gap = Inches(0.2)
    
    if element.direction == 'horizontal':
        panel_w = (target_w - (gap * (num_panels - 1))) / num_panels
        max_bottom = target_y
        for i, panel in enumerate(element.panels):
            px = target_x + i * (panel_w + gap)
            py = target_y
            if panel.title:
                tb = ctx.slide.shapes.add_textbox(px, py, panel_w, Inches(0.4))
                p = tb.text_frame.paragraphs[0]
                p.text = panel.title
                p.font.name = ctx.theme.font_name
                p.font.size = ctx.theme.size_body
                p.font.bold = True
                p.font.color.rgb = ctx.theme.color_primary
                py += Inches(0.5)
            end_y = py
            for pe in panel.elements:
                pe_ph = ctx.find_placeholder(getattr(pe, 'placeholder', None))
                if pe_ph:
                    render_element(pe, ctx, pe_ph.left, pe_ph.top, pe_ph.width, pe_ph.height)
                else:
                    end_y = render_element(pe, ctx, px, end_y, panel_w, target_h - (py - target_y))
            if end_y > max_bottom:
                max_bottom = end_y
        if not ph:
            rendered_height = getattr(element, 'height_hint', None)
            if rendered_height is None:
                rendered_height = max_bottom - target_y
            return y + rendered_height + ELEMENT_GAP
    else:
        panel_w = target_w
        panel_h = (target_h - (gap * (num_panels - 1))) / num_panels
        py = target_y
        for i, panel in enumerate(element.panels):
            px = target_x
            panel_start_y = py
            if panel.title:
                tb = ctx.slide.shapes.add_textbox(px, py, panel_w, Inches(0.4))
                p = tb.text_frame.paragraphs[0]
                p.text = panel.title
                p.font.name = ctx.theme.font_name
                p.font.size = ctx.theme.size_body
                p.font.bold = True
                p.font.color.rgb = ctx.theme.color_primary
                py += Inches(0.5)
            for idx_pe, pe in enumerate(panel.elements):
                pe_ph = ctx.find_placeholder(getattr(pe, 'placeholder', None))
                if pe_ph:
                    render_element(pe, ctx, pe_ph.left, pe_ph.top, pe_ph.width, pe_ph.height)
                else:
                    adj_h = get_adjusted_height(panel.elements, idx_pe, panel_start_y + panel_h, py, panel_w)
                    py = render_element(pe, ctx, px, py, panel_w, adj_h)
            py = panel_start_y + panel_h + gap
        if not ph:
            rendered_height = getattr(element, 'height_hint', None)
            if rendered_height is None:
                rendered_height = py - target_y - gap
            return y + rendered_height + ELEMENT_GAP
    return y
