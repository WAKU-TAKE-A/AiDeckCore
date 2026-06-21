from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from ..render_context import SlideContext
from ..height_estimator import ELEMENT_GAP

def render(element, ctx: SlideContext, x, y, w, h) -> float:
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    start_x = ph.left if ph else x
    start_y = ph.top if ph else y
    
    flat_nodes = []
    def build_flat_list(node, level, parent_idx):
        current_idx = len(flat_nodes)
        flat_nodes.append((node, level, parent_idx))
        for child in node.children:
            build_flat_list(child, level + 1, current_idx)
    build_flat_list(element.root, 0, -1)
    
    node_width = Inches(1.5)
    node_height = Inches(0.4)
    vertical_gap = Inches(0.15)
    horizontal_gap = Inches(0.4)
    indent_width = node_width + horizontal_gap
    
    from collections import defaultdict
    children_by_parent = defaultdict(list)
    for i, (node, level, parent_idx) in enumerate(flat_nodes):
        if parent_idx != -1:
            children_by_parent[parent_idx].append(i)
    
    leaf_counter = 0
    node_y = {}
    def calculate_y(idx):
        nonlocal leaf_counter
        children = children_by_parent[idx]
        if not children:
            ny = start_y + leaf_counter * (node_height + vertical_gap)
            node_y[idx] = ny
            leaf_counter += 1
            return ny
        else:
            child_ys = [calculate_y(c_idx) for c_idx in children]
            ny = sum(child_ys) / len(child_ys)
            node_y[idx] = ny
            return ny
    calculate_y(0)
    
    node_shapes = {}
    for i, (node, level, parent_idx) in enumerate(flat_nodes):
        nx = start_x + level * indent_width
        ny = node_y[i]
        
        shape = ctx.slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, nx, ny, node_width, node_height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = ctx.theme.color_flow_fill
        shape.line.color.rgb = ctx.theme.color_flow_line
        shape.line.width = Pt(1.5)
        
        tf = shape.text_frame
        tf.word_wrap = True
        tf.clear()
        p = tf.paragraphs[0]
        p.text = node.label
        p.font.name = ctx.theme.font_name
        p.font.size = ctx.theme.size_body_small
        p.font.color.rgb = ctx.theme.color_flow_text
        p.alignment = PP_ALIGN.CENTER
        
        node_shapes[i] = shape
    
    for parent_idx, child_indices in children_by_parent.items():
        for c_idx in child_indices:
            connector = ctx.slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, 0, 0, 0, 0)
            connector.line.color.rgb = ctx.theme.color_flow_line
            connector.line.width = Pt(1.5)
            connector.begin_connect(node_shapes[parent_idx], 3)
            connector.end_connect(node_shapes[c_idx], 1)
            
    if not ph:
        rendered_height = getattr(element, 'height_hint', None)
        if rendered_height is None:
            leaf_count = leaf_counter
            rendered_height = leaf_count * node_height + (leaf_count - 1) * vertical_gap
        return y + rendered_height + ELEMENT_GAP
    return y
