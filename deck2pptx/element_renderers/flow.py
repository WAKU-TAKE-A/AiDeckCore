from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from ..render_context import SlideContext
from ..height_estimator import ELEMENT_GAP

def render(element, ctx: SlideContext, x, y, w, h) -> float:
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    start_x = ph.left if ph else x
    start_y = ph.top if ph else y

    node_width = Inches(1.5)
    node_height = Inches(0.5)
    
    node_positions = {}
    
    def style_flow_node(shape):
        shape.fill.solid()
        shape.fill.fore_color.rgb = ctx.theme.color_flow_fill
        shape.line.color.rgb = ctx.theme.color_flow_line
        for p in shape.text_frame.paragraphs:
            p.font.name = ctx.theme.font_name
            p.font.size = ctx.theme.size_body_small
            p.font.color.rgb = ctx.theme.color_flow_text
    
    def style_flow_arrow(shape):
        shape.fill.solid()
        shape.fill.fore_color.rgb = ctx.theme.color_flow_line
        shape.line.color.rgb = ctx.theme.color_flow_line
    
    if element.direction == 'horizontal':
        for i, node in enumerate(element.nodes):
            nx = start_x + i * (node_width + Inches(0.5))
            ny = start_y
            shape = ctx.slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                nx, ny, node_width, node_height
            )
            shape.text = node.label
            style_flow_node(shape)
            node_positions[f"{node.id}_out"] = (nx + node_width, ny + node_height / 2)
            node_positions[f"{node.id}_in"] = (nx, ny + node_height / 2)
            
        for edge in element.edges:
            if f"{edge.from_node}_out" in node_positions and f"{edge.to_node}_in" in node_positions:
                fx, fy = node_positions[f"{edge.from_node}_out"]
                tx, ty = node_positions[f"{edge.to_node}_in"]
                
                if tx >= fx:
                    arrow_x = fx
                    arrow_width = max(tx - fx, Inches(0.1))
                    shape_type = MSO_SHAPE.RIGHT_ARROW
                else:
                    arrow_x = tx
                    arrow_width = max(fx - tx, Inches(0.1))
                    shape_type = MSO_SHAPE.LEFT_ARROW
                    
                arrow = ctx.slide.shapes.add_shape(
                    shape_type,
                    arrow_x, fy - Inches(0.1), arrow_width, Inches(0.2)
                )
                style_flow_arrow(arrow)
    else: # vertical
        for i, node in enumerate(element.nodes):
            nx = start_x
            ny = start_y + i * (node_height + Inches(0.5))
            shape = ctx.slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                nx, ny, node_width, node_height
            )
            shape.text = node.label
            style_flow_node(shape)
            node_positions[f"{node.id}_out"] = (nx + node_width / 2, ny + node_height)
            node_positions[f"{node.id}_in"] = (nx + node_width / 2, ny)
    
        for edge in element.edges:
            if f"{edge.from_node}_out" in node_positions and f"{edge.to_node}_in" in node_positions:
                fx, fy = node_positions[f"{edge.from_node}_out"]
                tx, ty = node_positions[f"{edge.to_node}_in"]
                if ty >= fy:
                    arrow_y = fy
                    arrow_height = max(ty - fy, Inches(0.1))
                    shape_type = MSO_SHAPE.DOWN_ARROW
                else:
                    arrow_y = ty
                    arrow_height = max(fy - ty, Inches(0.1))
                    shape_type = MSO_SHAPE.UP_ARROW
                    
                arrow = ctx.slide.shapes.add_shape(
                    shape_type,
                    fx - Inches(0.1), arrow_y, Inches(0.2), arrow_height
                )
                style_flow_arrow(arrow)

    if not ph:
        rendered_height = getattr(element, 'height_hint', None)
        if rendered_height is None:
            if element.direction == 'horizontal':
                rendered_height = node_height
            else:
                rendered_height = len(element.nodes) * (node_height + Inches(0.5))
        return y + rendered_height + ELEMENT_GAP
    return y
