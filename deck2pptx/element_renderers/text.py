from pptx.util import Pt
from ..render_context import SlideContext
from ..height_estimator import _estimate_element_height

def render(element, ctx: SlideContext, x, y, w, h) -> float:
    from ..render_utils import _set_text_frame_text
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    if ph and ph.has_text_frame:
        _set_text_frame_text(ph.text_frame, element.content)
        return y
    else:
        target_x = ph.left if ph else x
        target_y = ph.top if ph else y
        target_w = ph.width if ph else w
        rendered_height = _estimate_element_height(element, target_w, ctx.calibrated_metrics, ctx.theme, ctx.level_fonts, ctx.calibrated_heights)
        if hasattr(element, 'content') and '棚は万なり' in element.content:
            print(f"DEBUG: '棚は万なり' text.render rendered_height = {rendered_height / 914400.0} inches")
        txBox = ctx.slide.shapes.add_textbox(target_x, target_y, target_w, rendered_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = ctx.theme.text.padding
        tf.margin_bottom = ctx.theme.text.padding
        _set_text_frame_text(
            tf,
            element.content,
            font_name=ctx.theme.font.name,
            font_size=Pt(ctx.level_fonts[0]) if 0 in ctx.level_fonts else ctx.theme.font.size_body,
            font_color=ctx.theme.color.text,
        )
        if not ph:
            return y + rendered_height + ctx.theme.layout.element_gap
        return y

def render_bullet(element, ctx: SlideContext, x, y, w, h) -> float:
    from ..models import ListItem
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    if ph and ph.has_text_frame:
        tf = ph.text_frame
        tf.clear()
        for i, item in enumerate(element.items):
            is_li = isinstance(item, ListItem)
            text = item.text if is_li else str(item)
            level = item.level if is_li else 0

            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            if level > 0: p.level = level

            font_size = ctx.level_fonts.get(level)

            if font_size:
                p.font.size = Pt(font_size)
        return y
    else:
        target_x = ph.left if ph else x
        target_y = ph.top if ph else y
        target_w = ph.width if ph else w
        rendered_height = _estimate_element_height(element, target_w, ctx.calibrated_metrics, ctx.theme, ctx.level_fonts, ctx.calibrated_heights)
        txBox = ctx.slide.shapes.add_textbox(target_x, target_y, target_w, rendered_height)
        tf = txBox.text_frame
        tf.word_wrap = True
        tf.margin_top = ctx.theme.bullet.padding
        tf.margin_bottom = ctx.theme.bullet.padding
        for i, item in enumerate(element.items):
            is_li = isinstance(item, ListItem)
            text = item.text if is_li else str(item)
            level = item.level if is_li else 0

            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            if level > 0: p.level = level
            p.font.name = ctx.theme.font.name

            font_size = ctx.level_fonts.get(level)

            p.font.size = Pt(font_size) if font_size else (ctx.theme.font.size_body if level == 0 else ctx.theme.font.size_body_small)
            p.font.color.rgb = ctx.theme.color.text
        if not ph:
            return y + rendered_height + ctx.theme.layout.element_gap
        return y
