from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from ..render_context import SlideContext

def render(element, ctx: SlideContext, x, y, w, h) -> float:
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    start_x = ph.left if ph else x
    start_y = ph.top if ph else y
    width = ph.width if ph else w
    caption_height = ctx.theme.code.caption_height

    if element.caption or element.language:
        tb = ctx.slide.shapes.add_textbox(start_x, start_y, width, caption_height)
        p = tb.text_frame.paragraphs[0]
        p.text = element.caption if element.caption else f"Language: {element.language}"
        p.font.name = ctx.theme.font.name
        p.font.size = ctx.theme.font.size_body_semi_small
        p.font.italic = True
        start_y += caption_height

    line_count = len(element.code.splitlines()) if element.code else 1
    box_height = Inches(line_count * ctx.theme.code.line_height_factor + ctx.theme.code.height_padding)

    shape = ctx.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, start_x, start_y, width, box_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ctx.theme.color.surface
    shape.line.color.rgb = ctx.theme.color.border

    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = element.code
    p.font.name = ctx.theme.font.name_code
    p.font.size = ctx.theme.font.size_body_small
    p.font.color.rgb = ctx.theme.color.text
    p.alignment = PP_ALIGN.LEFT

    if not ph:
        rendered_height = getattr(element, 'height_hint', None)
        if rendered_height is None:
            rendered_height = (caption_height if element.caption or element.language else 0) + box_height
        return y + rendered_height + ctx.theme.layout.element_gap
    return y
