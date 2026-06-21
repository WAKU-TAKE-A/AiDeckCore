from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from ..render_context import SlideContext
from ..height_estimator import ELEMENT_GAP

def render(element, ctx: SlideContext, x, y, w, h) -> float:
    img_path = ctx.base_dir / element.source
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    try:
        if ph and hasattr(ph, 'insert_picture'):
            pic = ph.insert_picture(str(img_path))
            if element.caption:
                tb = ctx.slide.shapes.add_textbox(pic.left, pic.top + pic.height, pic.width, Inches(0.3))
                p = tb.text_frame.paragraphs[0]
                p.text = element.caption
                p.font.name = ctx.theme.font_name
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = ctx.theme.color_text_light
            return y
        else:
            target_x = ph.left if ph else x
            target_y = ph.top if ph else y
            max_w = ph.width if ph else w
            max_h = ph.height if ph else h
            if element.caption:
                max_h -= Inches(0.3)
                
            from PIL import Image as PILImage
            try:
                with PILImage.open(str(img_path)) as pil_img:
                    img_w, img_h = pil_img.size
                ratio = min(max_w / img_w, max_h / img_h)
                new_w = img_w * ratio
                new_h = img_h * ratio
            except:
                new_w = max_w
                new_h = None
                
            pic = ctx.slide.shapes.add_picture(str(img_path), target_x, target_y, width=new_w, height=new_h)
            if element.caption:
                tb = ctx.slide.shapes.add_textbox(target_x, target_y + (new_h if new_h else pic.height), new_w, Inches(0.3))
                p = tb.text_frame.paragraphs[0]
                p.text = element.caption
                p.font.name = ctx.theme.font_name
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = ctx.theme.color_text_light
                
            if not ph:
                rendered_height = getattr(element, 'height_hint', None)
                if rendered_height is None:
                    rendered_height = (new_h if new_h else pic.height) + (Inches(0.3) if element.caption else 0)
                return y + rendered_height + ELEMENT_GAP
            return y
    except Exception as e:
        print(f"Failed to load image {img_path}: {e}")
        return y

def render_gallery(element, ctx: SlideContext, x, y, w, h) -> float:
    num_images = len(element.images)
    if num_images == 0:
        return y
        
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    cols = getattr(element, 'columns', None)
    rows_ct = getattr(element, 'rows', None)
    
    if cols is None and rows_ct is None:
        if num_images == 1: cols, rows_ct = 1, 1
        elif num_images == 2: cols, rows_ct = 2, 1
        elif num_images == 3: cols, rows_ct = 3, 1
        elif num_images == 4: cols, rows_ct = 2, 2
        else: cols, rows_ct = 3, 2
    elif cols is None:
        import math
        cols = math.ceil(num_images / rows_ct)
    elif rows_ct is None:
        import math
        rows_ct = math.ceil(num_images / cols)
        
    cell_width = (ph.width if ph else w) / cols
    cell_height = (ph.height if ph else h) / rows_ct
    
    for i, img in enumerate(element.images[:cols*rows_ct]):
        r = i // cols
        c = i % cols
        cell_x = (ph.left if ph else x) + (c * cell_width)
        cell_y = (ph.top if ph else y) + (r * cell_height)
        
        max_w = cell_width - Inches(0.1)
        max_h = cell_height - Inches(0.1)
        if img.caption:
            max_h -= Inches(0.3)
            
        img_path = ctx.base_dir / img.source
        try:
            from PIL import Image as PILImage
            with PILImage.open(str(img_path)) as pil_img:
                img_w, img_h = pil_img.size
            ratio = min(max_w / img_w, max_h / img_h)
            new_w = img_w * ratio
            new_h = img_h * ratio
            
            center_x = cell_x + (cell_width - new_w) / 2
            center_y = cell_y + (cell_height - new_h - (Inches(0.3) if img.caption else 0)) / 2
            
            ctx.slide.shapes.add_picture(str(img_path), center_x, center_y, width=new_w, height=new_h)
            
            if img.caption:
                tb = ctx.slide.shapes.add_textbox(cell_x, center_y + new_h, cell_width, Inches(0.3))
                p = tb.text_frame.paragraphs[0]
                p.text = img.caption
                p.font.name = ctx.theme.font_name
                p.font.size = Pt(12)
                p.alignment = PP_ALIGN.CENTER
                p.font.color.rgb = ctx.theme.color_text_light
        except Exception as e:
            print(f"Failed to load image {img_path}: {e}")
    
    if not ph:
        rendered_height = getattr(element, 'height_hint', None)
        if rendered_height is None:
            rendered_height = cell_height * rows_ct
        return y + rendered_height + ELEMENT_GAP
    return y
