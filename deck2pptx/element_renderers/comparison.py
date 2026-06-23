from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from ..render_context import SlideContext
from ..height_estimator import ELEMENT_GAP
from ..text_utils import count_rendered_lines

def render(element, ctx: SlideContext, x, y, w, h) -> float:
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    num_cols = len(element.columns)
    if num_cols == 0:
        return y
    
    col_width = (ph.width if ph else w) / num_cols
    start_x = ph.left if ph else x
    start_y = ph.top if ph else y
    
    if element.title:
        tb = ctx.slide.shapes.add_textbox(start_x, start_y, col_width * num_cols, Inches(0.5))
        p = tb.text_frame.paragraphs[0]
        p.text = element.title
        p.font.name = ctx.theme.font_name
        p.font.size = ctx.theme.size_title
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER
        start_y += Inches(0.5)

    calib_h = None
    calib_cpi = None
    if ctx.calibrated_metrics and ctx.theme.size_body in ctx.calibrated_metrics:
        calib_h = ctx.calibrated_metrics[ctx.theme.size_body]['line_height']
        calib_cpi = ctx.calibrated_metrics[ctx.theme.size_body]['chars_per_inch']
    
    col_width_inches = col_width / 914400.0
    chars_per_line = max(1, int(col_width_inches * calib_cpi)) if calib_cpi else 30
    line_height = calib_h if calib_h else Inches(0.35)

    max_lines = 0
    for col in element.columns:
        lines = 1
        for item in col.items:
            lines += count_rendered_lines(item, chars_per_line)
        if lines > max_lines:
            max_lines = lines
            
    box_height = max_lines * line_height + Inches(0.1)

    for i, col in enumerate(element.columns):
        col_x = start_x + (i * col_width)
        
        tb = ctx.slide.shapes.add_textbox(col_x, start_y, col_width, box_height)
        tf = tb.text_frame
        tf.word_wrap = True
        
        p = tf.paragraphs[0]
        p.text = col.label
        p.font.name = ctx.theme.font_name
        p.font.size = ctx.theme.size_body
        p.font.bold = True
        p.alignment = PP_ALIGN.LEFT
        
        for item in col.items:
            p = tf.add_paragraph()
            p.text = f"• {item}"
            p.level = 0
            p.font.name = ctx.theme.font_name
            p.font.size = ctx.theme.size_body_semi_small
            
    if not ph:
        rendered_height = getattr(element, 'height_hint', None)
        if rendered_height is None:
            rendered_height = (Inches(0.5) if element.title else 0) + box_height
        return y + rendered_height + ELEMENT_GAP
    return y

def render_timeline(element, ctx: SlideContext, x, y, w, h) -> float:
    ph = ctx.find_placeholder(getattr(element, 'placeholder', None))
    start_x = ph.left if ph else x
    start_y = ph.top if ph else y
    width = ph.width if ph else w
    
    event_height = Inches(0.8)
    for i, ev in enumerate(element.events):
        ev_y = start_y + (i * event_height)
        
        # Date/Title box
        tb = ctx.slide.shapes.add_textbox(start_x, ev_y, Inches(2.0), event_height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ev.label
        p.alignment = PP_ALIGN.RIGHT
        p.font.name = ctx.theme.font_name
        p.font.size = ctx.theme.size_title
        p.font.bold = True
        p.font.color.rgb = ctx.theme.color_flow_line
        
        # Vertical Line
        line = ctx.slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, start_x + Inches(2.1), ev_y + Inches(0.1), Inches(0.05), event_height - Inches(0.2))
        line.fill.solid()
        line.fill.fore_color.rgb = ctx.theme.color_flow_line
        line.line.color.rgb = ctx.theme.color_flow_line
        
        # Content box
        tb = ctx.slide.shapes.add_textbox(start_x + Inches(2.3), ev_y, width - Inches(2.3), event_height)
        tf = tb.text_frame
        p = tf.paragraphs[0]
        p.text = ev.title
        p.font.name = ctx.theme.font_name
        p.font.size = ctx.theme.size_body
        p.font.bold = True
        p.font.color.rgb = ctx.theme.color_flow_text
        
        if ev.description:
            p2 = tf.add_paragraph()
            p2.text = ev.description
            p2.font.name = ctx.theme.font_name
            p2.font.size = ctx.theme.size_body_semi_small
            p2.font.color.rgb = ctx.theme.color_text_light
            
    if not ph:
        rendered_height = getattr(element, 'height_hint', None)
        if rendered_height is None:
            rendered_height = len(element.events) * event_height
        return y + rendered_height + ELEMENT_GAP
    return y
