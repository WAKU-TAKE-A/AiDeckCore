from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import re
from .models import Deck, Slide, BulletList
from .layout import Layout, get_slide_layout_type
from .theme import Theme

from .render_context import SlideContext
from .height_estimator import get_adjusted_height
from .element_renderers import render_element

def _remove_existing_slides(prs):
    """Use a template for layouts/theme without carrying starter slides forward."""
    slide_id_list = prs.slides._sldIdLst
    for slide_id in list(slide_id_list):
        prs.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)

def _name_matches(actual_name: str, requested_name: str) -> bool:
    actual = actual_name.strip().casefold()
    requested = requested_name.strip().casefold()
    return bool(requested) and actual.startswith(requested)

def _name_equals(actual_name: str, requested_name: str) -> bool:
    return actual_name.strip().casefold() == requested_name.strip().casefold()



def render_deck(deck: Deck, output_path: str, base_dir: Path = Path('.'), template_path: str = None, calib_first_slide: bool = False):
    calibrated_metrics = {}
    title_metrics = None
    theme = Theme(deck.theme)

    # Initialize Presentation
    if template_path:
        prs = Presentation(template_path)

        # Calibration Extraction
        do_calib = calib_first_slide
        if not do_calib and len(prs.slides) > 0:
            calib_slide = prs.slides[0]
            if calib_slide.shapes.title and calib_slide.shapes.title.has_text_frame:
                title_text = calib_slide.shapes.title.text.strip().lower()
                if title_text == "calibration" or re.match(r"^キャリ.ブレーション$", title_text):
                    do_calib = True
                    
        if do_calib and len(prs.slides) > 0:
            calib_slide = prs.slides[0]
            from .height_estimator import extract_template_metrics
            calibrated_metrics, title_metrics = extract_template_metrics(calib_slide, theme)

        level_fonts = {}
        if calibrated_metrics:
            sorted_fonts = sorted(calibrated_metrics.keys(), reverse=True)
            for i, fs in enumerate(sorted_fonts):
                level_fonts[i] = fs

        _remove_existing_slides(prs)
    else:
        prs = Presentation()
        level_fonts = {}

    base_dir = Path(base_dir)
    from .height_estimator import calibrate_line_heights
    calibrated_heights = calibrate_line_heights(deck, theme)

    # Set slide dimensions (only when no template; templates keep their own size)
    if not template_path:
        if deck.orientation == 'portrait':
            prs.slide_width = theme.layout.slide_width_portrait
            prs.slide_height = theme.layout.slide_height_portrait
        else:
            prs.slide_width = theme.layout.slide_width_landscape
            prs.slide_height = theme.layout.slide_height_landscape
        
    layout = Layout(prs.slide_width, prs.slide_height, theme, title_metrics=title_metrics)

    render_slides = list(deck.slides)

    if deck.toc:
        toc_title_text = deck.toc_title or "Table of Contents"
        toc_items = []
        for i, s in enumerate(deck.slides):
            if i == 0 and get_slide_layout_type(s) == "title":
                continue
            if s.title:
                toc_items.append(s.title)

        toc_slide = Slide(
            title=toc_title_text,
            layout_hint="Title and Content",
            elements=[BulletList(items=toc_items)]
        )

        if len(render_slides) > 0:
            render_slides.insert(1, toc_slide)
        else:
            render_slides.append(toc_slide)

    for slide_model in render_slides:
        layout_type = get_slide_layout_type(slide_model)

        # Determine layout
        slide_layout = None
        if slide_model.layout_hint:
            # Try to find by name
            for ly in prs.slide_layouts:
                if _name_matches(ly.name, slide_model.layout_hint):
                    slide_layout = ly
                    break

        if not slide_layout:
            # fallback to blank (index 6 usually) or title (index 0)
            if layout_type == "title" and len(prs.slide_layouts) > 0:
                slide_layout = prs.slide_layouts[0]
            elif len(prs.slide_layouts) > 6:
                slide_layout = prs.slide_layouts[6]
            elif len(prs.slide_layouts) > 0:
                slide_layout = prs.slide_layouts[-1]

        slide = prs.slides.add_slide(slide_layout)

        # Add slide notes if any
        if slide_model.notes and slide.has_notes_slide:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_model.notes

        def find_placeholder(name):
            if not name: return None
            for shape in slide.shapes:
                if _name_equals(shape.name, name):
                    return shape
            matching_layout_idx = None
            for layout_shape in slide.slide_layout.placeholders:
                if _name_matches(layout_shape.name, name):
                    matching_layout_idx = layout_shape.placeholder_format.idx
                    break
            if matching_layout_idx is not None:
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == matching_layout_idx:
                        return shape
            for shape in slide.shapes:
                if _name_matches(shape.name, name):
                    return shape
            return None

        # Try to use default title/subtitle placeholders if no manual coords
        title_ph = find_placeholder("Title")
        subtitle_ph = find_placeholder("Subtitle")
        for shape in slide.shapes:
            if shape.is_placeholder:
                if not title_ph and shape.placeholder_format.type in (1, 3): # TITLE or CENTER_TITLE
                    title_ph = shape
                elif not subtitle_ph and shape.placeholder_format.type == 4: # SUBTITLE
                    subtitle_ph = shape

        title_slide_y = layout.slide_height * theme.layout.title_slide_position_ratio

        slide_no_ph = find_placeholder("slideno")
        if slide_no_ph:
            slide_no_ph.text = str(slide_idx + 1)
            
        section_no_ph = find_placeholder("sectionno")
        if section_no_ph and slide_model.section_no:
            section_no_ph.text = str(slide_model.section_no)

        footer_ph = find_placeholder("footer")
        if footer_ph and deck.footer is not None:
            footer_ph.text = str(deck.footer)

        date_ph = find_placeholder("date")
        if date_ph and deck.date is not None:
            date_ph.text = str(deck.date)
        if slide_model.title:
            if title_ph:
                title_ph.text = slide_model.title
            else:
                title_y = title_slide_y if layout_type == "title" else layout.title_y
                title_box_height = theme.layout.title_slide_title_box_height if layout_type == "title" else layout.title_height
                txBox = slide.shapes.add_textbox(layout.title_x, title_y, layout.title_width, title_box_height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_model.title
                if layout_type == "title": p.alignment = PP_ALIGN.CENTER
                p.font.name = theme.font.name
                p.font.size = theme.font.size_title
                p.font.color.rgb = theme.color.primary if layout_type != "title" else theme.color.text

        if slide_model.subtitle:
            if subtitle_ph:
                subtitle_ph.text = slide_model.subtitle
            else:
                # If we made a manual title box, add paragraph there if title existed, else make new box
                # For simplicity, fallback to adding a new textbox below title
                sub_y = (
                    (title_slide_y + theme.layout.title_slide_title_box_height)
                    if layout_type == "title" else (layout.title_y + layout.title_height)
                )
                txBox = slide.shapes.add_textbox(layout.title_x, sub_y, layout.title_width, theme.layout.title_slide_subtitle_box_height)
                tf = txBox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_model.subtitle
                if layout_type == "title": p.alignment = PP_ALIGN.CENTER
                p.font.name = theme.font.name
                p.font.size = theme.font.size_subtitle
                p.font.color.rgb = theme.color.text_light

        ctx = SlideContext(
            slide=slide,
            theme=theme,
            deck=deck,
            base_dir=base_dir,
            layout=layout,
            find_placeholder=find_placeholder,
            calibrated_metrics=calibrated_metrics,
            level_fonts=level_fonts,
            calibrated_heights=calibrated_heights
        )

        align_val = getattr(slide_model, 'content_align', None) or getattr(deck, 'content_align', None)
        y_offset = 0
        if align_val:
            val = align_val.lower()
            if val in ('top',): y_offset = -theme.layout.align_offset_large
            elif val in ('semi-top', 'high'): y_offset = -theme.layout.align_offset_small
            elif val in ('semi-bottom', 'low'): y_offset = theme.layout.align_offset_small
            elif val in ('bottom',): y_offset = theme.layout.align_offset_large

        current_y = layout.content_y + y_offset
        content_x = layout.content_x
        for idx_el, element in enumerate(slide_model.elements):
            ph = find_placeholder(getattr(element, "placeholder", None))
            adj_h = layout.content_height if ph else get_adjusted_height(
                slide_model.elements, idx_el,
                layout.content_y + layout.content_height,
                current_y, layout.content_width,
                calibrated_metrics=ctx.calibrated_metrics, theme=ctx.theme, level_fonts=ctx.level_fonts, calibrated_heights=ctx.calibrated_heights
            )
            current_y = render_element(element, ctx, content_x, current_y, layout.content_width, adj_h)

    output_path = Path(output_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
