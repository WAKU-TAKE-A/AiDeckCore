import subprocess
import tempfile
import shutil
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE

from .models import Text, BulletList, Image, Table, Gallery, Flow, Split, CodeBlock, Mermaid, Tree, Comparison, Timeline
from .text_utils import count_rendered_lines
from .theme import Theme

# Lazily-built default Theme, used only when a caller does not have one
# (e.g. isolated unit tests). The renderer always passes ctx.theme, so in
# the normal build pipeline this is never constructed.
_default_theme = None


def _theme_or_default(theme):
    global _default_theme
    if theme is not None:
        return theme
    if _default_theme is None:
        _default_theme = Theme()
    return _default_theme


def _get_font_size(level_fonts, theme, level=0):
    # NOTE: when `theme` is available (always, in the real renderer),
    # theme.font.size_body wins over the fallback below regardless of
    # `level` -- this mirrors the original (pre-theme.ini) behavior, where
    # the level-based fallback was effectively dead code once a Theme
    # existed. The fallback values are still configurable in theme.ini
    # for completeness, but tuning them will not change normal output.
    theme = _theme_or_default(theme)
    fs = theme.font.fallback_size_level0 if level == 0 else theme.font.fallback_size_other
    if theme and hasattr(theme, 'font') and hasattr(theme.font, 'size_body'):
        fs = theme.font.size_body.pt
    if level_fonts and level in level_fonts:
        fs = level_fonts[level]
    return fs


def _get_calibrated_metrics(fs, calibrated_metrics):
    if not calibrated_metrics:
        return None, None
    if fs in calibrated_metrics:
        return calibrated_metrics[fs]['height'], calibrated_metrics[fs]['cpi']

    # Find closest
    closest_fs = min(calibrated_metrics.keys(), key=lambda k: abs(k - fs))
    ref = calibrated_metrics[closest_fs]
    ratio = fs / closest_fs
    return ref['height'] * ratio, ref['cpi'] * (1.0 / ratio)

def _estimate_element_height(element, content_width, calibrated_metrics=None, theme=None, level_fonts=None, calibrated_heights=None):
    """Return estimated rendered height (EMU) for an element."""
    theme = _theme_or_default(theme)

    if getattr(element, 'height_hint', None) is not None:
        return element.height_hint

    if isinstance(element, Text):
        fs = _get_font_size(level_fonts, theme, 0)
        calib_h, calib_cpi = _get_calibrated_metrics(fs, calibrated_metrics)
        
        line_height = calib_h
        if not line_height and calibrated_heights:
            line_height = calibrated_heights.get(0)
            
        if line_height:
            if calib_cpi:
                width_inches = content_width / theme.calibration.emu_per_inch if content_width else theme.calibration.fallback_width_inches
                chars_per_line = max(1, int(width_inches * calib_cpi))
            else:
                chars_per_line = theme.text.fallback_chars_per_line
            estimated_lines = element.content.count('\n') + 1 + len(element.content) // chars_per_line
            return estimated_lines * line_height + theme.text.padding
        else:
            estimated_lines = element.content.count('\n') + 1 + len(element.content) // theme.text.fallback_chars_per_line
            return estimated_lines * theme.text.line_height

    if isinstance(element, BulletList):
        from .models import ListItem
        total_h = 0

        for item in element.items:
            lvl = item.level if isinstance(item, ListItem) else 0
            text = item.text if isinstance(item, ListItem) else str(item)
            fs = _get_font_size(level_fonts, theme, lvl)
            calib_h, calib_cpi = _get_calibrated_metrics(fs, calibrated_metrics)
            line_height = calib_h
            if not line_height and calibrated_heights:
                line_height = calibrated_heights.get(lvl)

            if line_height:
                if calib_cpi:
                    width_inches = content_width / theme.calibration.emu_per_inch if content_width else theme.calibration.fallback_width_inches
                    indent_inches = theme.bullet.indent_per_level * (lvl + 1)
                    avail_width = max(1.0, width_inches - indent_inches)
                    chars_per_line = max(1, int(avail_width * calib_cpi))
                else:
                    chars_per_line = theme.text.fallback_chars_per_line
                lines = text.count('\n') + 1 + len(text) // chars_per_line
                total_h += lines * line_height
            else:
                weight = theme.bullet.level_weight_default if lvl == 0 else theme.bullet.level_weight_indented
                total_h += weight * theme.bullet.line_height
        return total_h + theme.bullet.padding

    if isinstance(element, CodeBlock):
        lines = len(element.code.splitlines()) if element.code else 1
        box_h = Inches(lines * theme.code.line_height_factor + theme.code.height_padding)
        caption_h = theme.code.caption_height if (getattr(element, 'caption', None) or getattr(element, 'language', None)) else 0
        return caption_h + box_h

    if isinstance(element, Mermaid):
        from . import mermaid_handler
        if mermaid_handler.has_mermaid_cli():
            return getattr(element, 'height_hint', None) or theme.mermaid.default_height
        else:
            lines = len(element.code.splitlines()) if element.code else 1
            box_h = Inches(lines * theme.code.line_height_factor + theme.code.height_padding)
            return box_h

    if isinstance(element, Tree):
        def count_leaves(node):
            if not node.children:
                return 1
            return sum(count_leaves(child) for child in node.children)
        leaf_count = count_leaves(element.root)
        node_height = theme.tree.node_height
        vertical_gap = theme.tree.vertical_gap
        return max(theme.layout.min_remaining_height, leaf_count * node_height + (leaf_count - 1) * vertical_gap)

    if isinstance(element, Timeline):
        return len(element.events) * theme.timeline.event_height

    if isinstance(element, Comparison):
        num_cols = len(element.columns)
        if num_cols == 0:
            return theme.layout.min_remaining_height
        col_width_inches = (content_width / theme.calibration.emu_per_inch if content_width else theme.calibration.fallback_width_inches) / num_cols

        calib_h = None
        calib_cpi = None
        body_font_size = theme.font.size_body
        if calibrated_metrics and body_font_size in calibrated_metrics:
            calib_h = calibrated_metrics[body_font_size].get('height')
            calib_cpi = calibrated_metrics[body_font_size].get('cpi')
        
        chars_per_line = max(1, int(col_width_inches * calib_cpi)) if calib_cpi else theme.comparison.fallback_chars_per_line
        line_height = calib_h if calib_h else theme.text.line_height

        max_lines = 0
        for col in element.columns:
            lines = 1  # for the label
            for item in col.items:
                lines += count_rendered_lines(item, chars_per_line)
            if lines > max_lines:
                max_lines = lines

        title_h = theme.comparison.title_height if element.title else 0
        return title_h + (max_lines * line_height) + theme.comparison.padding

    if isinstance(element, Flow):
        if getattr(element, 'direction', 'horizontal') == 'vertical':
            count = len(element.items) if hasattr(element, 'items') else 1
            return max(Inches(1.0), count * Inches(0.5) + max(0, count - 1) * Inches(0.4))
        else:
            return Inches(1.5)

    # Image, Gallery, Table, Split:
    # height depends on runtime data — caller handles these separately.
    return theme.layout.min_remaining_height

def get_adjusted_height(elements_list, current_idx, total_bottom_y, current_y, content_width, calibrated_metrics=None, theme=None, level_fonts=None, calibrated_heights=None):
    """Estimate available height for elements_list[current_idx]."""
    theme = _theme_or_default(theme)

    remaining_h = total_bottom_y - current_y
    if remaining_h < theme.layout.min_remaining_height:
        remaining_h = theme.layout.min_remaining_height

    from . import mermaid_handler
    has_mmdc = mermaid_handler.has_mermaid_cli()

    # Reserve height for all text-like elements that come AFTER the current one.
    future_elements = elements_list[current_idx + 1:]
    reserved_text_h = sum(
        _estimate_element_height(e, content_width, calibrated_metrics, theme, level_fonts, calibrated_heights) + theme.layout.element_gap
        for e in future_elements
        if getattr(e, 'placeholder', None) is None and (
            isinstance(e, (Text, BulletList, CodeBlock, Tree, Comparison, Timeline)) or (isinstance(e, Mermaid) and not has_mmdc)
        )
    )

    remaining_imgs = sum(
        1 for e in elements_list[current_idx:]
        if isinstance(e, (Image, Gallery, Split, Table)) or (isinstance(e, Mermaid) and has_mmdc)
    )

    available_img_h = remaining_h - reserved_text_h
    if available_img_h < theme.layout.min_remaining_height:
        available_img_h = theme.layout.min_remaining_height

    current_element = elements_list[current_idx]
    if (isinstance(current_element, (Image, Gallery, Split, Table)) or (isinstance(current_element, Mermaid) and has_mmdc)) and remaining_imgs > 0:
        return available_img_h / remaining_imgs
    else:
        return remaining_h

def calibrate_line_heights(deck, theme):
    soffice_cmd = None
    if shutil.which("soffice"):
        soffice_cmd = "soffice"
    else:
        win_path = r"C:\Program Files\LibreOffice\program\soffice.exe"
        if Path(win_path).exists():
            soffice_cmd = win_path
    
    if not soffice_cmd:
        return {}

    prs = Presentation()
    blank_layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[0]
    slide = prs.slides.add_slide(blank_layout)

    default_sizes = [24, 20, 18, 16, 14]
    
    for level in range(5):
        fs = getattr(deck, f"font_size_l{level}", None)
        if fs is None:
            fs = default_sizes[level]
        
        txBox = slide.shapes.add_textbox(Inches(0), Inches(level * 1.5), Inches(20), Inches(1))
        tf = txBox.text_frame
        tf.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        tf.word_wrap = False
        
        lines_text = "\n".join([f"Line {i+1}" for i in range(10)])
        tf.text = lines_text
        
        for p in tf.paragraphs:
            p.font.name = theme.font.name if theme and hasattr(theme, 'font') and theme.font.name else "Arial"
            p.font.size = Pt(fs)

    temp_dir = tempfile.mkdtemp()
    try:
        input_pptx = Path(temp_dir) / "calib_input.pptx"
        prs.save(str(input_pptx))
        
        output_temp_dir = Path(tempfile.mkdtemp())
        try:
            subprocess.run([soffice_cmd, "--headless", "--convert-to", "pptx", str(input_pptx), "--outdir", str(output_temp_dir)], check=True, capture_output=True)
            output_pptx = output_temp_dir / "calib_input.pptx"
            if not output_pptx.exists():
                return {}
            
            prs_out = Presentation(str(output_pptx))
            slide_out = prs_out.slides[0]
            
            result = {}
            for level, shape in enumerate(slide_out.shapes):
                if shape.height == Inches(1):
                    return {}
                result[level] = shape.height / 10.0
                
            return result
        finally:
            shutil.rmtree(output_temp_dir, ignore_errors=True)
    except Exception as e:
        print(f"Calibration error: {e}")
        return {}
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)

def extract_template_metrics(slide, theme=None) -> tuple[dict, dict]:
    """Extracts calibration metrics from a slide."""
    calibrated_metrics = {}
    title_metrics = {
        'font_size_pt': 20.0,
        'height_inches': 1.2,
        'top_inches': 0.5,
        'left_inches': 0.5,
        'width_inches': 0.0
    }
    title_found = False
    
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format.type in (1, 3):
            if not title_found:
                title_metrics['height_inches'] = shape.height / 914400.0 if shape.height else 1.2
                title_metrics['top_inches'] = shape.top / 914400.0 if shape.top else 0.5
                title_metrics['left_inches'] = shape.left / 914400.0 if shape.left else 0.5
                title_metrics['width_inches'] = shape.width / 914400.0 if shape.width else 0.0
                
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        for run in p.runs:
                            if run.font.size:
                                title_metrics['font_size_pt'] = run.font.size.pt
                                break
                        if title_metrics['font_size_pt'] != 20.0: break
                title_found = True
            continue
            
        if shape.has_text_frame:
            text = shape.text
            lines = text.count('\n') + text.count('\x0b') + 1
            if lines >= 3:
                font_size_pt = None
                for p in shape.text_frame.paragraphs:
                    for run in p.runs:
                        if run.font.size:
                            font_size_pt = run.font.size.pt
                            break
                    if font_size_pt: break
                
                if font_size_pt:
                    if shape.height:
                        height_per_line = int(shape.height / lines)
                    else:
                        theme = _theme_or_default(theme)
                        lh_factor = theme.calibration.line_height_factor if theme else 1.2
                        height_per_line = int(font_size_pt * 12700 * lh_factor)
                        
                    first_para_text = shape.text_frame.paragraphs[0].text.split('\x0b')[0]
                    cpi = len(first_para_text) / (shape.width / 914400.0) if shape.width else 60.0 / 6.0
                    calibrated_metrics[font_size_pt] = {
                        'shape_name': shape.name,
                        'font_size_pt': font_size_pt,
                        'lines': lines,
                        'chars_per_line': len(first_para_text),
                        'box_width_inches': shape.width / 914400.0 if shape.width else 0.0,
                        'box_height_inches': shape.height / 914400.0 if shape.height else 0.0,
                        'cpi': cpi,
                        'height': height_per_line
                    }
    return calibrated_metrics, title_metrics

