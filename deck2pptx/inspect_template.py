from pptx import Presentation
from pathlib import Path
import json
import sys

def extract_template_info(pptx_path: str | Path) -> dict:
    prs = Presentation(pptx_path)
    layouts_info = []
    
    for idx, layout in enumerate(prs.slide_layouts):
        placeholders = []
        for shape in layout.placeholders:
            ph_format = shape.placeholder_format
            placeholders.append({
                "idx": ph_format.idx,
                "name": shape.name,
                "type": str(ph_format.type)
            })
            
        shapes = []
        for shape in layout.shapes:
            if not shape.is_placeholder:
                shapes.append({
                    "name": shape.name,
                    "type": type(shape).__name__
                })
                
        layouts_info.append({
            "index": idx,
            "name": layout.name,
            "placeholders": placeholders,
            "other_shapes": shapes
        })
        
    return {
        "template": str(pptx_path),
        "layouts": layouts_info
    }

def inspect_template(pptx_path: str | Path, output_format: str = 'text'):
    try:
        result = extract_template_info(pptx_path)
    except Exception as e:
        if output_format == 'json':
            print(json.dumps({"error": str(e)}, indent=2))
        else:
            print(f"Error loading template: {e}", file=sys.stderr)
        sys.exit(1)
    
    if output_format == 'json':
        print(json.dumps(result, indent=2))
    else:
        print(f"Template: {result['template']}\n")
        for layout in result['layouts']:
            print(f"Layout {layout['index']}: '{layout['name']}'")
            if layout['placeholders']:
                print("  Placeholders:")
                for ph in layout['placeholders']:
                    print(f"    - [{ph['idx']}] '{ph['name']}' (Type: {ph['type']})")
            if layout['other_shapes']:
                print("  Other Shapes:")
                for sh in layout['other_shapes']:
                    print(f"    - '{sh['name']}' ({sh['type']})")
            print()
