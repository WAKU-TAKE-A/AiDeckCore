from pptx import Presentation
from pptx.util import Inches, Length
prs = Presentation('Inputs/template_tate_01.pptx')
slide = prs.slides[0]
for i, shape in enumerate(slide.shapes):
    if shape.has_text_frame:
        lines = shape.text.count('\n') + 1
        h = Length(shape.height)
        h_per_line = Length(int(shape.height / lines))
        print(f'Shape {i}: height={h.inches:.2f}in, lines={lines}, height_per_line={h_per_line.inches:.2f}in')
