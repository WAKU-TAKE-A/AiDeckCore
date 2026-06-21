from pathlib import Path

import pytest
from PIL import Image as PILImage
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches

from deck2pptx.adapters import load_deck
from deck2pptx.layout import Layout
from deck2pptx.models import Deck, Gallery, Image, Slide
from deck2pptx.renderer import render_deck
from deck2pptx.validation import ValidationError, validate_deck


def _make_image(path: Path, size=(1200, 500), color=(180, 210, 230)):
    PILImage.new("RGB", size, color).save(path)


def test_markdown_image_alt_is_preserved_as_caption(tmp_path):
    _make_image(tmp_path / "photo.jpg")
    md = tmp_path / "deck.md"
    md.write_text("# Image\n\n![hoge06](photo.jpg)\n", encoding="utf-8")

    deck = load_deck(md)

    image = deck.slides[0].elements[0]
    assert isinstance(image, Image)
    assert image.source == "photo.jpg"
    assert image.caption == "hoge06"


def test_markdown_consecutive_images_preserve_gallery_captions(tmp_path):
    _make_image(tmp_path / "a.jpg")
    _make_image(tmp_path / "b.jpg")
    md = tmp_path / "deck.md"
    md.write_text("# Gallery\n\n<!-- gallery -->\n![first](a.jpg)\n![second](b.jpg)\n", encoding="utf-8")

    deck = load_deck(md)

    gallery = deck.slides[0].elements[0]
    assert isinstance(gallery, Gallery)
    assert len(gallery.images) == 2
    assert gallery.images[0].caption == "first"
    assert gallery.images[1].caption == "second"


def test_yaml_image_caption_and_gallery_grid_are_parsed(tmp_path):
    _make_image(tmp_path / "a.jpg")
    _make_image(tmp_path / "b.jpg")
    yml = tmp_path / "deck.yaml"
    yml.write_text(
        """
slides:
  - title: Images
    elements:
      - image:
          source: a.jpg
          caption: Single caption
      - gallery:
          rows: 2
          columns: 2
          images:
            - source: a.jpg
              caption: A caption
            - b.jpg
""",
        encoding="utf-8",
    )

    deck = load_deck(yml)
    single = deck.slides[0].elements[0]
    gallery = deck.slides[0].elements[1]

    assert isinstance(single, Image)
    assert single.caption == "Single caption"
    assert isinstance(gallery, Gallery)
    assert gallery.rows == 2
    assert gallery.columns == 2
    assert [img.source for img in gallery.images] == ["a.jpg", "b.jpg"]
    assert [img.caption for img in gallery.images] == ["A caption", None]


def test_yaml_legacy_gallery_image_strings_still_work(tmp_path):
    _make_image(tmp_path / "a.jpg")
    _make_image(tmp_path / "b.jpg")
    yml = tmp_path / "deck.yaml"
    yml.write_text(
        """
slides:
  - title: Legacy Gallery
    elements:
      - gallery:
          images:
            - a.jpg
            - b.jpg
""",
        encoding="utf-8",
    )

    deck = load_deck(yml)
    gallery = deck.slides[0].elements[0]

    assert isinstance(gallery, Gallery)
    assert gallery.rows is None
    assert gallery.columns is None
    assert [img.source for img in gallery.images] == ["a.jpg", "b.jpg"]


@pytest.mark.parametrize(
    ("rows", "columns", "expected"),
    [
        (0, 2, "invalid_gallery_rows"),
        (2, 0, "invalid_gallery_columns"),
        ("2", 2, "invalid_gallery_rows"),
        (2, "2", "invalid_gallery_columns"),
    ],
)
def test_invalid_gallery_grid_rejected(rows, columns, expected):
    deck = Deck(slides=[Slide(title="Bad", elements=[Gallery(images=[], rows=rows, columns=columns)])])

    with pytest.raises(ValidationError) as excinfo:
        validate_deck(deck, Path("."))

    assert expected in [error["code"] for error in excinfo.value.errors]


def test_gallery_rendering_fits_images_and_captions_inside_content_area(tmp_path):
    for idx in range(4):
        _make_image(tmp_path / f"wide{idx}.jpg", size=(1600, 450), color=(150 + idx * 20, 190, 220))

    deck = Deck(
        slides=[
            Slide(
                title="Gallery",
                elements=[
                    Gallery(
                        rows=2,
                        columns=2,
                        images=[
                            Image(source=f"wide{idx}.jpg", caption=f"hoge{idx}")
                            for idx in range(4)
                        ],
                    )
                ],
            )
        ]
    )
    out = tmp_path / "gallery.pptx"

    render_deck(deck, str(out), base_dir=tmp_path)

    prs = Presentation(str(out))
    slide = prs.slides[0]
    layout = Layout(prs.slide_width, prs.slide_height)
    pictures = [shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.PICTURE]
    captions = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text.startswith("hoge")
    ]

    assert len(pictures) == 4
    assert len(captions) == 4

    for picture in pictures:
        assert picture.left >= layout.content_x
        assert picture.top >= layout.content_y
        assert picture.left + picture.width <= layout.content_x + layout.content_width
        assert picture.top + picture.height <= layout.content_y + layout.content_height

    for caption in captions:
        assert caption.top + caption.height <= layout.content_y + layout.content_height
        assert caption.text_frame.paragraphs[0].font.size.pt == 12
