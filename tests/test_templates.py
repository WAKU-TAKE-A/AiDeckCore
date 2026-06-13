import pytest
from pathlib import Path
from deck2pptx.inspect_template import extract_template_info
from deck2pptx.renderer import render_deck
from deck2pptx.models import Deck, Slide, Text
from deck2pptx.yaml_adapter import load_yaml
from deck2pptx.markdown_adapter import load_markdown
from pptx import Presentation

def test_extract_template_info(tmp_path):
    # Generate a dummy template
    template_path = tmp_path / "dummy_template.pptx"
    prs = Presentation()
    # It has 11 default layouts. Let's just use it as is.
    prs.save(str(template_path))
    
    info = extract_template_info(str(template_path))
    assert info["template"] == str(template_path)
    assert len(info["layouts"]) > 0
    
    title_layout = info["layouts"][0]
    assert title_layout["index"] == 0
    assert "name" in title_layout

def test_yaml_parity():
    yaml_content = """
title: Test YAML Parity
toc: true
font_size_l0: 24
slides:
  - title: Slide 1
    layout_hint: Title Slide
    elements:
      - text: Hello
        placeholder: Text Placeholder 1
"""
    import yaml
    from deck2pptx.yaml_adapter import load_yaml
    import tempfile
    
    with tempfile.NamedTemporaryFile(mode='w', delete=False, suffix=".yaml", encoding='utf-8') as f:
        f.write(yaml_content)
        tmp_name = f.name
        
    try:
        deck = load_yaml(tmp_name)
        assert deck.toc is True
        assert deck.font_size_l0 == 24
        assert len(deck.slides) == 1
        assert deck.slides[0].layout_hint == "Title Slide"
        assert deck.slides[0].elements[0].placeholder == "Text Placeholder 1"
    finally:
        Path(tmp_name).unlink()

def test_markdown_new_page_and_layout():
    md_content = """
<!-- layout="Title Slide" -->
# Slide 1
<!-- subtitle="Sub" -->

<!-- new_page="Blank" -->
# Slide 2

<!-- new_page -->
Just some text without heading.
"""
    from deck2pptx.markdown_adapter import load_markdown
    import tempfile
    
    with tempfile.NamedTemporaryFile(mode='w', delete=False, suffix=".md", encoding='utf-8') as f:
        f.write(md_content)
        tmp_name = f.name
        
    try:
        deck = load_markdown(tmp_name)
        assert len(deck.slides) == 3
        assert deck.slides[0].title == "Slide 1"
        assert deck.slides[0].subtitle == "Sub"
        assert deck.slides[0].layout_hint == "Title Slide"
        
        assert deck.slides[1].title == "Slide 2"
        assert deck.slides[1].layout_hint == "Blank"
        
        # The new_page without heading
        assert deck.slides[2].title == ""
        assert len(deck.slides[2].elements) == 1
        assert deck.slides[2].elements[0].content == "Just some text without heading."
    finally:
        Path(tmp_name).unlink()

def test_build_with_template(tmp_path):
    # Generate a dummy template
    template_path = tmp_path / "dummy_template.pptx"
    prs = Presentation()
    prs.save(str(template_path))
    
    # Create a deck that uses a layout from the dummy template
    # The default blank template usually has "Title Slide" as layouts[0]
    deck = Deck(title="Test", orientation="landscape", theme="default")
    slide1 = Slide(title="Slide 1", layout_hint="Title Slide")
    # This element uses a placeholder that probably doesn't exist, testing the fallback
    slide1.elements.append(Text(content="Fallback text", placeholder="Nonexistent Placeholder"))
    deck.slides.append(slide1)
    
    output_path = tmp_path / "output.pptx"
    
    # Should not crash, fallback should trigger
    render_deck(deck, str(output_path), base_dir=tmp_path, template_path=str(template_path))
    
    assert output_path.exists()
    
    # Verify the created presentation
    out_prs = Presentation(str(output_path))
    assert len(out_prs.slides) == 1
    assert out_prs.slides[0].shapes.title.text == "Slide 1"

def test_layout_lookup_uses_case_insensitive_prefix(tmp_path):
    template_path = tmp_path / "layout_prefix_template.pptx"
    prs = Presentation()
    prs.save(str(template_path))

    deck = Deck(title="Test", orientation="landscape", theme="default")
    deck.slides.append(Slide(title="Prefix Layout", layout_hint="title sl"))

    output_path = tmp_path / "output.pptx"
    render_deck(deck, str(output_path), base_dir=tmp_path, template_path=str(template_path))

    out_prs = Presentation(str(output_path))

    assert len(out_prs.slides) == 1
    assert out_prs.slides[0].slide_layout.name == "Title Slide"
    assert out_prs.slides[0].shapes.title.text == "Prefix Layout"

def test_build_with_template_does_not_copy_existing_template_slides(tmp_path):
    template_path = tmp_path / "template_with_starter_slide.pptx"
    prs = Presentation()
    starter = prs.slides.add_slide(prs.slide_layouts[0])
    starter.shapes.title.text = "Starter slide should not be copied"
    prs.save(str(template_path))

    deck = Deck(title="Test", orientation="landscape", theme="default")
    deck.slides.append(Slide(title="Generated Slide", layout_hint="Title Slide"))

    output_path = tmp_path / "output.pptx"
    render_deck(deck, str(output_path), base_dir=tmp_path, template_path=str(template_path))

    out_prs = Presentation(str(output_path))
    texts = [
        shape.text
        for slide in out_prs.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text
    ]

    assert len(out_prs.slides) == 1
    assert "Generated Slide" in texts
    assert "Starter slide should not be copied" not in texts

def test_placeholder_lookup_accepts_powerpoint_number_suffix(tmp_path):
    template_path = tmp_path / "numbered_placeholder_template.pptx"
    prs = Presentation()
    prs.save(str(template_path))

    deck = Deck(title="Test", orientation="landscape", theme="default")
    slide = Slide(title="Title", layout_hint="Title Slide")
    slide.elements.append(Text(content="Subtitle by prefix", placeholder="Subtitle"))
    deck.slides.append(slide)

    output_path = tmp_path / "output.pptx"
    render_deck(deck, str(output_path), base_dir=tmp_path, template_path=str(template_path))

    out_prs = Presentation(str(output_path))
    subtitle_shapes = [
        shape
        for shape in out_prs.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
        and shape.name.casefold().startswith("subtitle")
    ]

    assert subtitle_shapes
    assert subtitle_shapes[0].text == "Subtitle by prefix"

def test_placeholder_lookup_uses_case_insensitive_prefix(tmp_path):
    template_path = tmp_path / "placeholder_prefix_template.pptx"
    prs = Presentation()
    prs.save(str(template_path))

    deck = Deck(title="Test", orientation="landscape", theme="default")
    slide = Slide(title="Title", layout_hint="Title Slide")
    slide.elements.append(Text(content="Subtitle by lower prefix", placeholder="sub"))
    deck.slides.append(slide)

    output_path = tmp_path / "output.pptx"
    render_deck(deck, str(output_path), base_dir=tmp_path, template_path=str(template_path))

    out_prs = Presentation(str(output_path))
    subtitle_shapes = [
        shape
        for shape in out_prs.slides[0].shapes
        if getattr(shape, "has_text_frame", False)
        and shape.name.casefold().startswith("subtitle")
    ]

    assert subtitle_shapes
    assert subtitle_shapes[0].text == "Subtitle by lower prefix"

def test_placeholder_lookup_uses_layout_placeholder_name_when_slide_name_changes(tmp_path):
    template_path = tmp_path / "renamed_layout_placeholder_template.pptx"
    prs = Presentation()
    layout = prs.slide_layouts[1]
    body_idx = None
    for placeholder in layout.placeholders:
        if placeholder.is_placeholder and placeholder.placeholder_format.type == 7:
            placeholder.name = "Body"
            body_idx = placeholder.placeholder_format.idx
            break
    assert body_idx is not None
    prs.save(str(template_path))

    deck = Deck(title="Test", orientation="landscape", theme="default")
    slide = Slide(title="Title", layout_hint=layout.name)
    slide.elements.append(Text(content="Body by layout name", placeholder="Body"))
    deck.slides.append(slide)

    output_path = tmp_path / "output.pptx"
    render_deck(deck, str(output_path), base_dir=tmp_path, template_path=str(template_path))

    out_prs = Presentation(str(output_path))
    body_placeholder = out_prs.slides[0].placeholders[body_idx]

    assert body_placeholder.text == "Body by layout name"

def test_title_prefers_placeholder_named_title_over_title_type(tmp_path):
    template_path = tmp_path / "custom_title_placeholder_template.pptx"
    prs = Presentation()
    layout = prs.slide_layouts[0]
    named_title_idx = None
    title_type_idx = None
    for placeholder in layout.placeholders:
        if placeholder.placeholder_format.type == 3:
            placeholder.name = "ToCompany"
            title_type_idx = placeholder.placeholder_format.idx
        elif placeholder.placeholder_format.type == 4:
            placeholder.name = "Title"
            named_title_idx = placeholder.placeholder_format.idx
    assert named_title_idx is not None
    assert title_type_idx is not None
    prs.save(str(template_path))

    deck = Deck(title="Test", orientation="landscape", theme="default")
    deck.slides.append(Slide(title="Generated Title", layout_hint=layout.name))

    output_path = tmp_path / "output.pptx"
    render_deck(deck, str(output_path), base_dir=tmp_path, template_path=str(template_path))

    out_prs = Presentation(str(output_path))

    assert out_prs.slides[0].placeholders[named_title_idx].text == "Generated Title"
    assert out_prs.slides[0].placeholders[title_type_idx].text != "Generated Title"

def test_toc_rendering(tmp_path):
    # Create a deck with toc=True
    deck = Deck(title="Test", orientation="landscape", theme="default", toc=True, toc_title="Agenda")
    slide1 = Slide(title="Slide 1")
    slide2 = Slide(title="Slide 2")
    deck.slides.extend([slide1, slide2])
    
    output_path = tmp_path / "toc_output.pptx"
    render_deck(deck, str(output_path), base_dir=tmp_path)
    
    assert output_path.exists()
    out_prs = Presentation(str(output_path))
    
    # Should have 3 slides (Slide 1, Agenda, Slide 2) because render_deck injects TOC at index 1
    assert len(out_prs.slides) == 3
    assert out_prs.slides[1].shapes.title.text == "Agenda"
    
    # The TOC slide should contain a BulletList with the titles of the other slides
    # We can check if 'Slide 1' and 'Slide 2' are in the text (but not the title shape)
    text_content = ""
    for shape in out_prs.slides[1].shapes:
        if shape.has_text_frame and shape != out_prs.slides[1].shapes.title:
            text_content += shape.text_frame.text
    
    assert "Slide 1" in text_content
    assert "Slide 2" in text_content

def test_renderer_purity(tmp_path):
    # Regression test: render_deck must not mutate deck.slides when inserting TOC
    deck = Deck(title="Test Purity", orientation="landscape", theme="default", toc=True, toc_title="Agenda")
    deck.slides.extend([Slide(title="Slide 1"), Slide(title="Slide 2")])
    
    initial_slides_count = len(deck.slides)
    
    out1 = tmp_path / "out1.pptx"
    out2 = tmp_path / "out2.pptx"
    
    render_deck(deck, str(out1), base_dir=tmp_path)
    # deck.slides should be unchanged
    assert len(deck.slides) == initial_slides_count
    
    render_deck(deck, str(out2), base_dir=tmp_path)
    assert len(deck.slides) == initial_slides_count
    
    prs1 = Presentation(str(out1))
    prs2 = Presentation(str(out2))
    
    # Rendered output should have the exact same number of slides (initial + 1 TOC)
    assert len(prs1.slides) == initial_slides_count + 1
    assert len(prs2.slides) == initial_slides_count + 1

def test_list_hierarchy_rendering(tmp_path):
    from deck2pptx.models import ListItem, BulletList
    deck = Deck(title="Test", font_size_l0=24, font_size_l1=20, indent=20)
    slide1 = Slide(title="Slide 1")
    items = [
        ListItem(text="Level 0 item", level=0),
        ListItem(text="Level 1 item", level=1)
    ]
    slide1.elements.append(BulletList(items=items))
    deck.slides.append(slide1)
    
    output_path = tmp_path / "list_output.pptx"
    render_deck(deck, str(output_path), base_dir=tmp_path)
    
    out_prs = Presentation(str(output_path))
    # Check paragraphs in the slide
    found_text = False
    for shape in out_prs.slides[0].shapes:
        if shape.has_text_frame:
            tf = shape.text_frame
            for p in tf.paragraphs:
                if p.text == "Level 1 item":
                    found_text = True
                    assert p.level == 1
                    assert p.font.size.pt == 20
    assert found_text
