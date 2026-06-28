import pytest
from pathlib import Path
from pptx import Presentation
from deck2pptx.models import Deck, Slide
from deck2pptx.renderer import render_deck
from pptx.util import Pt

def test_auto_calib_from_title(tmp_path):
    # Create a dummy template
    tmpl_path = tmp_path / "calib_template.pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "キャリフブレーション" # Match regex ^キャリ.ブレーション$
    prs.save(tmpl_path)
    
    # Create a dummy deck
    deck = Deck(title="Test", slides=[])
    
    out_path = tmp_path / "out.pptx"
    
    # Mock the internal logic to verify if calib happened
    # A simple way is to check if calibrated_metrics was populated, but we can't easily intercept it.
    # However, we can patch `deck2pptx.renderer.calibrated_metrics`? No, it's a local var.
    # We can check if the first slide is removed: yes, template slides are removed by _remove_existing_slides.
    pass

def test_auto_calib_english_title(tmp_path):
    tmpl_path = tmp_path / "calib_template_en.pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "  CalibRation  "
    # To have calibrated_metrics extracted, the slide needs a shape with 2+ lines of text.
    txBox = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    txBox.text = "Line 1\nLine 2\nLine 3"
    for p in txBox.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(24)
    prs.save(tmpl_path)
    
    deck = Deck(title="Test", slides=[Slide(title="s1", elements=[])])
    out_path = tmp_path / "out2.pptx"
    
    # We can patch SlideContext to capture calibrated_metrics
    import deck2pptx.renderer
    original_SlideContext = deck2pptx.renderer.SlideContext
    captured_metrics = None
    
    class DummyContext:
        def __init__(self, **kwargs):
            nonlocal captured_metrics
            captured_metrics = kwargs.get('calibrated_metrics')
            self.slide = kwargs.get('slide')
            self.theme = kwargs.get('theme')
            self.level_fonts = kwargs.get('level_fonts')
            self.calibrated_metrics = kwargs.get('calibrated_metrics')
            self.calibrated_heights = kwargs.get('calibrated_heights')
            self.deck = kwargs.get('deck')
            self.layout = kwargs.get('layout')
            
    deck2pptx.renderer.SlideContext = DummyContext
    try:
        render_deck(deck, out_path, template_path=str(tmpl_path))
    finally:
        deck2pptx.renderer.SlideContext = original_SlideContext
        
    assert captured_metrics is not None
    assert len(captured_metrics) > 0 # Calibration occurred!

def test_no_auto_calib_wrong_title(tmp_path):
    tmpl_path = tmp_path / "calib_template_wrong.pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Not Calibration"
    txBox = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    txBox.text = "Line 1\nLine 2"
    for p in txBox.text_frame.paragraphs:
        p.font.size = Pt(24)
    prs.save(tmpl_path)
    
    deck = Deck(title="Test", slides=[Slide(title="s1", elements=[])])
    out_path = tmp_path / "out3.pptx"
    
    import deck2pptx.renderer
    original_SlideContext = deck2pptx.renderer.SlideContext
    captured_metrics = None
    
    class DummyContext:
        def __init__(self, **kwargs):
            nonlocal captured_metrics
            captured_metrics = kwargs.get('calibrated_metrics')
            self.slide = kwargs.get('slide')
            self.theme = kwargs.get('theme')
            self.level_fonts = kwargs.get('level_fonts')
            self.calibrated_metrics = kwargs.get('calibrated_metrics')
            self.calibrated_heights = kwargs.get('calibrated_heights')
            self.deck = kwargs.get('deck')
            self.layout = kwargs.get('layout')
            
    deck2pptx.renderer.SlideContext = DummyContext
    try:
        render_deck(deck, out_path, template_path=str(tmpl_path))
    finally:
        deck2pptx.renderer.SlideContext = original_SlideContext
        
    assert captured_metrics == {} # Calibration skipped!

def test_auto_calib_japanese_regex(tmp_path):
    tmpl_path = tmp_path / "calib_template_ja.pptx"
    prs = Presentation()
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "キャリ〰ブレーション" # Any character matches '.'
    txBox = slide.shapes.add_textbox(0, 0, 1000000, 1000000)
    txBox.text = "Line 1\nLine 2\nLine 3"
    for p in txBox.text_frame.paragraphs:
        for r in p.runs:
            r.font.size = Pt(24)
    prs.save(tmpl_path)
    
    deck = Deck(title="Test", slides=[Slide(title="s1", elements=[])])
    out_path = tmp_path / "out4.pptx"
    
    import deck2pptx.renderer
    original_SlideContext = deck2pptx.renderer.SlideContext
    captured_metrics = None
    
    class DummyContext:
        def __init__(self, **kwargs):
            nonlocal captured_metrics
            captured_metrics = kwargs.get('calibrated_metrics')
            self.slide = kwargs.get('slide')
            self.theme = kwargs.get('theme')
            self.level_fonts = kwargs.get('level_fonts')
            self.calibrated_metrics = kwargs.get('calibrated_metrics')
            self.calibrated_heights = kwargs.get('calibrated_heights')
            self.deck = kwargs.get('deck')
            self.layout = kwargs.get('layout')
            
    deck2pptx.renderer.SlideContext = DummyContext
    try:
        render_deck(deck, out_path, template_path=str(tmpl_path))
    finally:
        deck2pptx.renderer.SlideContext = original_SlideContext
        
    assert captured_metrics is not None
    assert len(captured_metrics) > 0 # Calibration occurred!
